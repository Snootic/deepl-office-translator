from io import BytesIO
from docx import Document
from pptx import Presentation
import pandas as pd
from typing import Union, IO
import tiktoken
import os
import datetime
import json
import sys

class File:
    def __init__(self) -> None:
        pass
    
    def load_document(self, file: Union[str, IO[bytes]] = None):
        if isinstance(file, str):
            ext = os.path.splitext(file)[1].lower()
        else:
            ext = os.path.splitext(file.name)[1].lower()

        if ext == ".docx":
            return self.load_word(file)
        elif ext == ".xlsx":
            return self.load_excel(file)
        elif ext == ".pptx":
            return self.load_pptx(file)
        else:
            raise ValueError("Tipo de arquivo não suportado. Por favor, insira um arquivo .docx, .xlsx ou .pptx.")

    def load_excel(self, file: str | IO[bytes ]= None):
        pass
    
    def load_pptx(self, file: str | IO[bytes ]= None):
        try:
            prs = Presentation(file)
        except Exception as e:
            raise(e)

        properties_dir = dir(prs.core_properties)

        properties_keys = [k for k in properties_dir if not k.startswith("_")]
        
        properties = {}
        for key in properties_keys:
            properties[key] = prs.core_properties.__getattribute__(key)

        properties["slide count"] = len(prs.slides)

        tokens = []
        final_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if isinstance(shape.text, (int,float)) or not shape.text.strip():
                                        continue
                    final_text += shape.text
                    pretokens = self.tokenize(shape.text)
                    for token in pretokens:
                        tokens.append(token)

                if shape.has_table:
                        table = shape.table 
                        table_data = []
                        for row in table.rows:
                            for cell in row.cells:
                                if isinstance(cell, (int,float)) or not cell.text.strip():
                                        continue
                                final_text += cell.text.rstrip()
                                pretokens = self.tokenize(cell.text.rstrip())
                                for token in pretokens:
                                    tokens.append(token)
                if shape.has_chart:
                        try:
                            blob_stream = BytesIO(shape.chart._workbook.xlsx_part.blob)
                            sheets = pd.ExcelFile(blob_stream).sheet_names
                            df = pd.read_excel(blob_stream)
                            columns = df.columns.tolist()
                            
                            for index, column in enumerate(columns):
                                if isinstance(column, (int,float)) or not column.strip():
                                        continue
                                final_text += column
                                
                                pretokens = self.tokenize(column)
                                for token in pretokens:
                                    tokens.append(token)
                                
                            row_data = df.values.tolist()
                            for row in row_data:
                                for index, data in enumerate(row):
                                    if isinstance(data, (int,float)) or not data.strip():
                                        continue
                                    
                                    final_text += data
                                
                                    pretokens = self.tokenize(column)
                                    for token in pretokens:
                                        tokens.append(token)
                            
                            new_df = pd.DataFrame(data=row_data,columns=columns)
                            
                            blob = BytesIO()
                            new_df.to_excel(blob,index=False, sheet_name=sheets[0], engine='xlsxwriter')
                            blob.seek(0)
                            
                            blob_data = blob.getvalue()
                            
                            shape.chart._workbook.xlsx_part.blob = blob_data
                            
                        except Exception as e:
                            continue
                                
        properties["word count"] = len(final_text.split(" "))

        properties["tokens"] = tokens
        
        try:
            properties["created"] = datetime.datetime.isoformat(properties["created"])
            properties["last_printed"] = datetime.datetime.isoformat(properties["last_printed"])
            properties["modified"] = datetime.datetime.isoformat(properties["modified"])
        except:
            pass
        
        properties["tokens count"] = len(tokens)

        result = json.dumps(obj=properties,skipkeys=True, default=lambda o: '<not serializable>',indent=2,ensure_ascii=False)
        
        return result

    def load_word(self, file: str | IO[bytes ]= None):
        try:
            file = Document(file)
        except Exception as e:
            raise(e)

        properties_dir = dir(file.core_properties)
        
        properties_keys = [k for k in properties_dir if not k.startswith("_")]
        
        properties = {}
        for key in properties_keys:
            properties[key] = file.core_properties.__getattribute__(key)

        paragraphs_text = ""
        tokens = []
        for paragraph in file.paragraphs:
            text = paragraph.text.rstrip()
            text = text.lstrip()
            paragraphs_text += " "+text
            pretokens = self.tokenize(paragraph.text)
            for token in pretokens:
                tokens.append(token)

        words = paragraphs_text.split(" ")

        while "" in words:
            words.remove("")

        properties["word_count"] = len(words)

        properties["tokens"] = tokens
        
        
        try:
            properties["created"] = datetime.datetime.isoformat(properties["created"])
            properties["last_printed"] = datetime.datetime.isoformat(properties["last_printed"])
            properties["modified"] = datetime.datetime.isoformat(properties["modified"])
        except:
            pass
        properties["tokens_count"] = len(tokens)
             
        result = json.dumps(obj=properties,skipkeys=True, default=lambda o: '<not serializable>',indent=2,ensure_ascii=False)

        return result
    
    def tokenize(self,text):
        enc = tiktoken.encoding_for_model("gpt-4o")
        tokens = enc.encode(text)

        return tokens
    

file = File()