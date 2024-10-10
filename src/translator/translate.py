import os
import deepl
import pandas as pd
from docx import Document
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from openai import OpenAI
from io import BytesIO
from pypdf import PdfReader, PdfWriter
from sys import _getframe

class Translate:
    # All returns must be json parsable
    
    def __init__(self) -> None:
        self.model:str
        self.API_KEY: str
    
    def main(self,api_key: str, model: str, **kwargs):
        self.API_KEY = api_key
        self.model = model
        if model.__contains__("gpt"):
            self.client = OpenAI(api_key=api_key, **kwargs)
        else:
            self.translator = deepl.Translator(api_key)
   
    def _attribute_are_equals(self, attr1, attr2) -> bool:
        if attr1 == attr2:
            True
        if attr1 is None and attr2 == False:
            True
        if attr2 is None and attr1 == False:
            True
        return False
    
    def _runs_are_equals(self, run1, run2, **kwargs) -> bool:
        caller = _getframe(1).f_code.co_name
        
        if caller == "translate_word_preserve_format":
            font1 = run1
            font2 = run2
        else:
            font1 = run1.font
            font2 = run2.font
        
        if not self._attribute_are_equals(run1.font.size, run2.font.size):
            return False
        
        elif (hasattr(run1.font.color, 'rgb') and hasattr(run2.font.color, 'rgb')):
            if run1.font.color.rgb != run2.font.color.rgb:
                    return False
                
        elif 'came_from_pdf' in kwargs:
            if kwargs['came_from_pdf'] == False:
                if run1.font.name != run2.font.name:
                    return False
            
        elif not self._attribute_are_equals(font1.bold, font2.bold):
            return False
        
        elif not self._attribute_are_equals(font1.underline, font2.underline):
            return False
        
        elif not self._attribute_are_equals(font1.italic, font2.italic):
            return False
        
        return True
    
    def _process_runs(self, paragraph) -> None:
        runs = []
        previous_run = None
        for run in paragraph.runs:
            if run.text.strip() == "":
                continue
            if previous_run is not None and self._runs_are_equals(previous_run, run):
                if previous_run in runs:
                    id = runs.index(previous_run)
                    runs.remove(previous_run)
                    previous_run.text += run.text
                    runs.insert(id, previous_run)
                else:
                    previous_run.text += run.text
                    runs.append(previous_run)
            else:
                previous_run = run
                runs.append(previous_run)
                
        return runs
    
    def _manage_runs_text(self, runs):
        text_to_translate = ""
        if len(runs) == 0:
            return None
        
        elif len(runs) > 1:
            for run in runs:
                if run != runs[-1]:
                    text_to_translate += f"{run.text}"+"{end-run}"
                    continue
                text_to_translate += f"{run.text}"
        else:
            text_to_translate += runs[0].text

        if text_to_translate.strip() == "" or text_to_translate.isdigit():
            return None

        return text_to_translate
    
    def translate_document(self, file: str, output_path: str, target_lang: str, source_lang: str = None, **kwargs):
        if isinstance(file, str):
            ext = os.path.splitext(file)[1].lower()
        else:
            ext = os.path.splitext(file.name)[1].lower()

        if ext == ".docx":
            return self.translate_word_preserve_format(file, output_path, target_lang, source_lang, **kwargs)
        elif ext == ".xlsx":
            return self.translate_excel(file, output_path, target_lang, source_lang, **kwargs)
        elif ext == ".pptx":
            return self.translate_pptx(file, output_path, target_lang, source_lang, **kwargs)
        else:
            raise ValueError("Tipo de arquivo não suportado. Por favor, insira um arquivo .docx, .xlsx ou .pptx.")
    
    def translate_pptx(self, presentation_path: str, output_path: str, target_lang: str, source_lang: str = None, pages: int | list = None, glossary: deepl.GlossaryInfo | str = None, **kwargs):
        prs = Presentation(presentation_path)

        def translate_paragraph(shape):
            for paragraph in shape.text_frame.paragraphs:
                if not paragraph.text or paragraph.text.strip() == "":
                    continue
                
                runs = self._process_runs(paragraph)

                paragraph.clear()

                text_to_translate = self._manage_runs_text(runs)
                
                if text_to_translate is None:
                    continue
                
                in_bank = word_bank.get(text_to_translate, None)

                if in_bank is not None:
                    translated_text = word_bank[text_to_translate]
                else:
                    translated_text = self.translate_text(text_to_translate, target_lang, source_lang, **kwargs)
                
                runs_text = translated_text.split("{end-run}")

                for id, run in enumerate(runs):
                    try:
                        runs_text[id]
                    except IndexError:
                        continue

                    word_bank[run.text] = runs_text[id]

                    run_copy = run

                    run_copy.text = runs_text[id]

                    new_run = paragraph.add_run()

                    new_run.font.bold = run_copy.font.bold
                    new_run.font.italic = run_copy.font.italic
                    new_run.font.underline = run_copy.font.underline
                    new_run.font.size = run_copy.font.size
                    if hasattr(run.font.color, 'rgb'):
                        new_run.font.color.rgb = run_copy.font.color.rgb
                    new_run.font.name = run_copy.font.name
                    new_run.text = run_copy.text

        def translate_table(shape):
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    original_text = cell.text.rstrip()
                    
                    if not original_text.strip() or original_text.isdigit():
                        continue
                    
                    if word_bank.get(original_text):
                        cell.text = word_bank[original_text]
                        continue
                    
                    translated_text = self.translate_text(original_text, target_lang, source_lang, **kwargs)
                    
                    # translated_text = "isso é uma célula de tabela"
                    
                    word_bank[original_text] = translated_text
                        
                    cell.text = translated_text

        def translate_chart(shape):
            blob_stream = BytesIO(shape.chart._workbook.xlsx_part.blob)
            # sheets = pd.ExcelFile(blob_stream).sheet_names
            df = pd.read_excel(blob_stream)
            blob_stream.close()
            
            columns = df.columns.tolist()
            
            for index, column in enumerate(columns):
                if not column.strip() or column.isdigit():
                    continue
                
                if word_bank.get(column):
                    columns[index] = word_bank[column]
                    continue
                
                new_column = self.translate_text(column, target_lang, source_lang, **kwargs)
                
                # new_column = "isso é uma coluna"
                
                word_bank[column] = new_column
                
                columns[index] = new_column
            columns.pop(0)
                
            row_data = df.values.tolist()
            categories = []
            for row in row_data:
                for index, data in enumerate(row):
                    if isinstance(data, (int,float)) or not data.strip():
                        continue
                    
                    if word_bank.get(data):
                        row[index] = word_bank[data]
                        continue
                    
                    row[index] = self.translate_text(data, target_lang, source_lang, **kwargs)
                    
                    # row[index] = "isso é uma linha"
                    
                    word_bank[data] = row[index]
                    
                categories.append(row[0])
                row.pop(0)
                
            df = None
            
            # new_df = pd.DataFrame(data=row_data,columns=columns)
            
            # blob = BytesIO()
            # new_df.to_excel(blob,index=False, sheet_name=sheets[0], engine='openpyxl')
            # blob.seek(0)
            
            # blob_data = blob.getvalue()
            
            chart_data = CategoryChartData()
            chart_data.categories = categories
            for index, column in enumerate(columns):
                series_data = tuple(row[index] for row in row_data)

                chart_data.add_series(column, series_data)

            shape.chart.replace_data(chart_data)

            # shape.chart._workbook.xlsx_part.blob = blob_data
            
        word_bank = {}
        try:
            if not pages:
                pages = [id + 1 for id, slide in enumerate(prs.slides)]

            for id, slide in enumerate(prs.slides):
                if isinstance(pages, int) and id != pages:
                    continue
                elif isinstance(pages, list) and id not in pages:
                    continue
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        translate_paragraph(shape)
                                    
                    if shape.has_table:
                        translate_table(shape)
                        
                    if shape.has_chart:
                        try:
                            translate_chart(shape)
                        except Exception as e:
                            print(e)            
        except:
            prs.save(output_path)

        prs.save(output_path)

    def deepl_translate_doc(self,doc_path: str, target_lang: str, source_lang: str = None, glossary: deepl.GlossaryInfo = None, **kwargs):
        if self.model.__contains__("gpt"):
            raise(TypeError("model not compatible with this type of translation"))
        output_path_list = doc_path.split('.')
        doc_extension = output_path_list[-1]
        output_path_list.pop()
        output_path_list.append(f"Translated {target_lang}")
        output_path_list.append(doc_extension)

        output_path = output_path_list[0]
        output_path_list.pop(0)

        for item in output_path_list:
            output_path += f".{item}"
        try:
            with open(doc_path, "rb") as in_file, open(output_path, "wb") as out_file:
                self.translator.translate_document(
                    in_file,
                    out_file,
                    source_lang=source_lang,
                    target_lang=target_lang,
                    glossary=glossary,
                    **kwargs
                )

        except deepl.DocumentTranslationException as error:
            doc_id = error.document_handle.document_id
            doc_key = error.document_handle.document_key
            print(f"Error after uploading ${error}, id: ${doc_id} key: ${doc_key}")
        except deepl.DeepLException as error:
            raise error

    def translate_word_preserve_format(self,doc_path: str, output_path: str, target_lang: str, source_lang: str = None, glossary: deepl.GlossaryInfo | str = None, **kwargs):
        doc = Document(doc_path)

        word_bank = {}

        def translate_paragraph(element):
            for paragraph in element.paragraphs:
                if not paragraph.text or paragraph.text.strip() =="":
                    continue
                
                runs = self._process_runs(paragraph)

                paragraph.clear()

                text_to_translate = self._manage_runs_text(runs)
                            
                if text_to_translate is None:
                    continue

                in_bank = word_bank.get(text_to_translate, None)

                if in_bank is not None:
                    translated_text = word_bank[text_to_translate]
                else:
                    translated_text = self.translate_text(text_to_translate, target_lang, source_lang, **kwargs)
                
                runs_text = translated_text.split("{end-run}")

                for id, run in enumerate(runs):
                    try:
                        print(runs_text[id])
                    except IndexError:
                        continue
                    word_bank[run.text] = runs_text[id]
                        
                    run.text = runs_text[id]

                    new_run = paragraph.add_run()

                    new_run.style = run.style
                    new_run.bold = run.bold
                    new_run.italic = run.italic
                    new_run.underline = run.underline
                    new_run.font.size = run.font.size
                    new_run.font.color.rgb = run.font.color.rgb
                    new_run.font.name = run.font.name

        def translate_tables():
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        translate_paragraph(cell)

        try:
            translate_paragraph(doc)
            translate_tables()
        except Exception as e:
            doc.save(output_path)
            return {"error": f"error but saved: {e}"}
        
        doc.save(output_path)

    def translate_excel(self, spreadsheet_path: str, output_path: str, source_column:str, target_column:str, target_lang: str,  source_lang: str = None, glossary: deepl.GlossaryInfo = None, **kwargs):
        df = pd.read_excel(spreadsheet_path)
        
        df[target_column] = df[source_column].apply(lambda x: self.translate_text(x, source_lang=source_lang, target_lang=target_lang, glossary=glossary, **kwargs))

        df.to_excel(output_path, index=False)  
    
    def translate_pdf(self, pdf_path: str, output_path: str, target_lang: str, source_lang: str = None, to_word: bool = False, **kwargs):
        reader = PdfReader(pdf_path)
        writer = PdfWriter(reader)
        
        for page in writer.pages:
            text = page.extract_text()
            
            # translated_text = self.translate_text(text, source_lang=source_lang, target_lang=target_lang, **kwargs)
            
            print(page.items)
            print("alo amigos")
           
    def deepl_translate_text(self, text: str, source_lang: str = None, target_lang: str = None, glossary: deepl.GlossaryInfo = None, **kwargs):
        try:
            result = self.translator.translate_text(text, source_lang=source_lang, target_lang=target_lang, glossary=glossary, **kwargs)
            return result.text
        except deepl.DeepLException as d:
            raise d
        except:
            return text

    def gpt_translate_text(self, text: str, target_language:str, source_language:str = None, context: str | None = None):
        prompt = f"Translate to {target_language}, return original text if you can't. Keep the formatting, capitalization, punctuation, use the fluent vocabulary of a native." +" Don't translate {end-run}, don't remove it."
        
        if context:
            prompt += " "+context
        completion = self.client.chat.completions.create(
            model=self.model,
            
            messages=[
                {
                    "role": "system", 
                    "content": prompt
                },
                {
                    "role": "user",
                    "content": text
                }
            ]
        )
        
        return completion.choices[0].message.content
    
    def translate_text(self, *args, **kwargs):
        if self.model.__contains__("gpt"):
            return self.gpt_translate_text(*args,**kwargs)
        return self.deepl_translate_text(*args,**kwargs)

translate = Translate()